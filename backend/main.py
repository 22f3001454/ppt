from fastapi import FastAPI, Form, HTTPException
from fastapi.responses import FileResponse
from fastapi.middleware.cors import CORSMiddleware
from pptx import Presentation
from pptx.util import Pt, Inches
import tempfile, os, json

app = FastAPI(title="Auto PPT Generator")
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

def local_fallback_split(text: str, tone: str):
    paras = [p.strip() for p in text.split("\n\n") if p.strip()]
    slides = []
    i = 0
    while i < len(paras) and len(slides) < 10:
        title = paras[i].split("\n")[0][:60]
        bullets = []
        for j in range(i, min(i + 2, len(paras))):
            sents = [s.strip() for s in paras[j].split(".") if s.strip()]
            bullets.extend(sents[:3])
        slides.append({"title": title or f"Slide {len(slides)+1}", "bullets": bullets})
        i += 2
    return slides

@app.post("/generate")
async def generate(
    text: str = Form(...),
    tone: str = Form(""),
    provider: str = Form("local"),
    api_key: str = Form(""),
):
    if not text.strip():
        raise HTTPException(status_code=400, detail="No input text")

    slides = local_fallback_split(text, tone)

    tmpdir = tempfile.mkdtemp()
    out_path = os.path.join(tmpdir, "generated.pptx")
    prs = Presentation()

    for s in slides:
        layout = prs.slide_layouts[1]  # Title and content
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = s["title"]

        body = slide.placeholders[1]
        tf = body.text_frame
        tf.clear()
        for bullet in s["bullets"]:
            p = tf.add_paragraph()
            p.text = bullet
            p.font.size = Pt(16)

    prs.save(out_path)

    preview = [{"title": s["title"], "bullets": s["bullets"][:3]} for s in slides[:5]]
    headers = {"x-slide-preview": json.dumps(preview)}

    return FileResponse(
        out_path,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename="generated_presentation.pptx",
        headers=headers,
    )
